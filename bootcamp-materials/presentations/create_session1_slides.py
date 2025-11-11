#!/usr/bin/env python3
"""
Create PowerPoint slides for Session 1: Prompt Files and Configuration
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_session1_slides():
    """Create Session 1 presentation slides"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SESSION 1: PROMPT FILES AND CONFIGURATION"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Create Your First Prompt Files for Spring Migration"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Duration
    duration_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    duration_frame = duration_box.text_frame
    duration_frame.text = "90 minutes"
    duration_para = duration_frame.paragraphs[0]
    duration_para.font.size = Pt(18)
    duration_para.alignment = PP_ALIGN.CENTER

    # Slide 2: From Individual Techniques to Organized Modes (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "From Individual Techniques to Organized Modes"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "You already know the 4 core techniques:"

    p = tf.add_paragraph()
    p.text = "Chain of Thought - Forces reasoning"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ReAct - Iterative problem-solving"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Tree of Thoughts - Explores alternatives"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Structured Prompts - Consistent output"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "The Problem:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "How do you ORGANIZE them?"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "How do you make them REUSABLE?"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Today's Answer: Prompt Files"
    p.font.bold = True
    p.font.size = Pt(20)

    # Slide 3: The 5 Prompt Modes Quick Overview (4 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The 5 Prompt Modes Quick Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Mode 1: Constitution - Project rules everywhere"

    p = tf.add_paragraph()
    p.text = "Example: javax → jakarta rules"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Mode 2: Specification - WHAT to build in THIS file"

    p = tf.add_paragraph()
    p.text = "Example: Migrate UserController specifics"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Mode 3: Planning - HOW to execute step-by-step"

    p = tf.add_paragraph()
    p.text = "Mode 4: ABCD - Handle decision points"

    p = tf.add_paragraph()
    p.text = "Mode 5: Implementation - Combine all and generate"

    # Slide 4: Spring Migration Example (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Spring Migration - Before & After"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = 'Bad Prompt: "Migrate this controller to Spring 3"'
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    p = tf.add_paragraph()
    p.text = "Result: Generic changes, might miss patterns"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = 'Better: Prompt Files Approach'
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Constitution file (reusable): javax → jakarta rules"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Specification file (per controller): UserController details"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Result: Systematic, reusable, documented"
    p.font.bold = True

    # Slide 5: Time Investment vs Payoff (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Time Investment vs Payoff"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Manual coding: 45 min every time"

    p = tf.add_paragraph()
    p.text = "With Prompt Files: 30 min first time"

    p = tf.add_paragraph()
    p.text = "Second file: 5 min (reuse Constitution)"
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Third+ files: 5 min each"
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Compound effect:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "After 5 migrations: 3+ hours saved"
    p.level = 1
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Slide 6: Today's Hands-On Plan (2 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Today's Hands-On Plan (60 minutes)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Setup Repository (5 min)"

    p = tf.add_paragraph()
    p.text = "Verify spring-migration-demo ready"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2. Create Constitution File (15 min)"

    p = tf.add_paragraph()
    p.text = "Spring 2→3 migration rules"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "3. Create Specification File (15 min)"

    p = tf.add_paragraph()
    p.text = "UserController migration details"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "4. Test with AI (20 min)"

    p = tf.add_paragraph()
    p.text = "5. Validate Output (10 min)"

    # Save presentation
    output_file = "SESSION-1-Prompt-Files.pptx"
    prs.save(output_file)
    print(f"✅ Created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    print("🎨 Creating Session 1 PowerPoint slides...")
    output = create_session1_slides()
    print(f"\n🎉 Done! Open {output} to view slides")
