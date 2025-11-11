#!/usr/bin/env python3
"""
Create PowerPoint slides for Session 1: Industry Standards & Practical Application
Updated with proper grounding in research
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_session1_slides_v2():
    """Create Session 1 presentation slides with industry standards framing"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SESSION 1: PROMPT ENGINEERING INDUSTRY STANDARDS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "From Ad-Hoc AI Prompting to Structured, Reusable Workflows"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    duration_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    duration_frame = duration_box.text_frame
    duration_frame.text = "90 minutes (15 min presentation + 60 min hands-on + 15 min review)"
    duration_para = duration_frame.paragraphs[0]
    duration_para.font.size = Pt(16)
    duration_para.alignment = PP_ALIGN.CENTER

    # Slide 2: The Problem - Ad-Hoc AI Prompting (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Problem: Ad-Hoc AI Prompting"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = 'Current State: "Hey AI, migrate this controller to Spring Boot 3"'
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Problems with Ad-Hoc Approach:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Every prompt starts from scratch"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "No shared team knowledge in prompts"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Inconsistent results across files"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Can't track what worked vs what didn't"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Becomes 'vibe coding' (unreviewed AI output)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Solution: Structured, versioned, reusable prompt patterns"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Slide 3: Industry Standards - The Three Tiers (4 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Industry Standards: The Three Tiers"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Tier 1: Proven Patterns (10+ years)"

    p = tf.add_paragraph()
    p.text = "Few-shot, Chain-of-Thought, Persona, Template patterns"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "ADRs (Architecture Decision Records), RFC Process"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Tier 2: Emerging Standards (1-3 years)"

    p = tf.add_paragraph()
    p.text = ".cursorrules, .github/copilot-instructions.md"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "3,000+ forks of awesome-cursorrules repository"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Tier 3: Experimental (months to 1-2 years)"

    p = tf.add_paragraph()
    p.text = "GitHub Spec-Kit (Sept 2024, ~14 months old as of Nov 2025)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Interesting experiments, growing understanding, not yet proven at scale"
    p.level = 1
    p.font.size = Pt(18)

    # Slide 4: Configuration Standards (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Configuration Standards: .cursorrules & Copilot Instructions"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Natural language files telling AI tools how to work with your project"

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = ".cursorrules Example:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Tech Stack: Spring Boot 3.2+, Java 17+, Jakarta EE"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Standards: DRY, KISS, YAGNI, constructor injection"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "DO NOT: Use deprecated javax.* packages"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Benefits:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Team-shared knowledge (version controlled)"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Consistent AI behavior across developers"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Reusable across projects"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Slide 5: Decision Documentation - ADRs vs Structured Prompts (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Decision Documentation: ADRs vs Structured Prompts"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Option A: Architecture Decision Records (Industry Standard)"

    p = tf.add_paragraph()
    p.text = "Proven since 2011, used by Microsoft, AWS, Google, Spotify"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Format: docs/adr/0001-use-jakarta.md"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Documents WHY decisions were made"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Option B: Structured Prompt Files (Experimental Pattern)"

    p = tf.add_paragraph()
    p.text = "AI-friendly format with rules and success criteria"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Focuses on HOW for AI generation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Both Are Valid! Can use BOTH together"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "ADRs document decisions, structured prompts guide AI"
    p.level = 1

    # Slide 6: Applying Standards - One Practical Example (2 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Applying Standards: One Practical Example"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Today's Hands-On: 5-File Workflow Pattern"

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "File 1: System Prompt (maps to: .cursorrules, persona pattern)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "File 2: Task Specification (maps to: template pattern)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "File 3: Execution Plan (maps to: task decomposition, ReAct)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "File 4: Decision Documentation (maps to: ADR, alternatives)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "File 5: Code Generation (maps to: few-shot + chain-of-thought)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Insight: This is ONE way to structure prompts"
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "ADRs alone, Spec-Kit, or custom patterns are equally valid"
    p.level = 1
    p.font.size = Pt(18)

    # Save presentation
    output_file = "SESSION-1-Industry-Standards.pptx"
    prs.save(output_file)
    print(f"✅ Created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    print("🎨 Creating Session 1 PowerPoint slides (Industry Standards version)...")
    output = create_session1_slides_v2()
    print(f"\n🎉 Done! Open {output} to view slides")
