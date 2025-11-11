#!/usr/bin/env python3
"""
Create PowerPoint slides for Session 2: Workflows and Orchestration
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_session2_slides():
    """Create Session 2 presentation slides"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SESSION 2: WORKFLOWS AND ORCHESTRATION"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Orchestrate All 5 Prompt Modes for Complete Migration Workflow"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Duration
    duration_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    duration_frame = duration_box.text_frame
    duration_frame.text = "90 minutes"
    duration_para = duration_frame.paragraphs[0]
    duration_para.font.size = Pt(18)
    duration_para.alignment = PP_ALIGN.CENTER

    # Slide 2: From Individual Modes to Complete Workflows (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "From Individual Modes to Complete Workflows"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Session 1 Recap:"

    p = tf.add_paragraph()
    p.text = "Constitution (Mode 1) - Reusable rules"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Specification (Mode 2) - File-specific targets"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Today's Goal:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Add Planning (Mode 3) - Execution steps"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Add ABCD (Mode 4) - Decision documentation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Add Implementation (Mode 5) - Synthesis and generation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Execute complete 5-mode workflow"
    p.level = 1
    p.font.bold = True

    # Slide 3: The 3 Missing Modes Quick Overview (4 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The 3 Missing Modes Quick Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Mode 3: Planning - Break work into ordered steps"

    p = tf.add_paragraph()
    p.text = "Phase 1: Imports (2 min), Phase 2: Annotations (5 min), Phase 3: Testing (5 min)"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Mode 4: ABCD Clarification - Document decisions"

    p = tf.add_paragraph()
    p.text = "Exception handling: Option A (simple) vs B (ProblemDetail) vs C (RFC 7807)"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Mode 5: Implementation - Pull everything together"

    p = tf.add_paragraph()
    p.text = "Context loaded: Constitution ✓, Specification ✓, Planning ✓, ABCD ✓ → Generate code"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Key Insight: Modes 1-4 provide context, Mode 5 triggers generation"
    p.font.bold = True
    p.font.size = Pt(20)

    # Slide 4: How Planning Mode Structures Execution (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How Planning Mode Structures Execution"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = 'Without Planning: AI makes random changes, unclear order'
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = 'With Planning:'
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Phase 1: Package Updates (2 min) - Must compile first"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Phase 2: Annotations (5 min) - Build on updated imports"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Phase 3: Testing (5 min) - Validate everything works"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Benefits: Logical order, small tasks, validation checkpoints, time estimates"
    p.font.bold = True

    # Slide 5: How ABCD Mode Documents Decisions (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How ABCD Mode Documents Decisions"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "When to Use ABCD:"

    p = tf.add_paragraph()
    p.text = "Multiple valid approaches exist"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Tradeoffs matter for your context"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Team needs to understand why you chose X"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Example: Exception Handling Approach"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "A: Simple String messages"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "B: ProblemDetail with details (recommended)"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "C: Full RFC 7807 (over-engineered)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "D: Keep current pattern"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Benefits: Documents WHY, team understands tradeoffs"
    p.font.bold = True

    # Slide 6: Complete 5-Mode Workflow (2 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Complete 5-Mode Workflow - Today's Plan"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Step 1: Create Planning File (15 min)"

    p = tf.add_paragraph()
    p.text = "Break UserController migration into phases"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Step 2: Create ABCD File (10 min)"

    p = tf.add_paragraph()
    p.text = "Document exception handling decision"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Step 3: Create Implementation File (10 min)"

    p = tf.add_paragraph()
    p.text = "Reference all 5 modes"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Step 4: Execute Workflow (15 min)"

    p = tf.add_paragraph()
    p.text = "Load all modes into AI, generate code"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Step 5: Validate & Document (10 min)"

    p = tf.add_paragraph()
    p.text = "Compare with feature branch"
    p.level = 1

    # Save presentation
    output_file = "SESSION-2-Workflows.pptx"
    prs.save(output_file)
    print(f"✅ Created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    print("🎨 Creating Session 2 PowerPoint slides...")
    output = create_session2_slides()
    print(f"\n🎉 Done! Open {output} to view slides")
