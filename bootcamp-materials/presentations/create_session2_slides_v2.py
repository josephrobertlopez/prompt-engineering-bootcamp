#!/usr/bin/env python3
"""
Create PowerPoint slides for Session 2: Advanced Patterns & Complete Workflows
Updated with proper grounding in research
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_session2_slides_v2():
    """Create Session 2 presentation slides with advanced patterns framing"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SESSION 2: ADVANCED PATTERNS & WORKFLOWS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Orchestrating Multiple Prompt Patterns for Complex Tasks"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    duration_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    duration_frame = duration_box.text_frame
    duration_frame.text = "90 minutes (15 min presentation + 60 min hands-on + 15 min review)"
    duration_para = duration_frame.paragraphs[0]
    duration_para.font.size = Pt(16)
    duration_para.alignment = PP_ALIGN.CENTER

    # Slide 2: From Individual Patterns to Orchestrated Workflows (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "From Individual Patterns to Orchestrated Workflows"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Session 1 Recap:"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Foundational patterns: Few-shot, Chain-of-Thought, Persona, Template"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Configuration standards: .cursorrules, Copilot instructions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Decision documentation: ADRs vs structured prompts"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Today's Goal:"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Learn advanced patterns: ReAct, Tree of Thoughts, Meta-prompting"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Orchestrate multiple patterns for complex tasks"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Compare workflow strategies at scale"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Execute complete end-to-end migration example"
    p.level = 1

    # Slide 3: Advanced Prompt Patterns Overview (4 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Advanced Prompt Patterns Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "ReAct Pattern (Reason + Act) - Yao et al. 2022"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Think → Act → Observe → Think → Act → Observe..."
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = "For multi-step tasks requiring validation at each stage"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Tree of Thoughts Pattern - Yao et al. 2023"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Generate multiple solution branches → Evaluate → Choose best"
    p.level = 1
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = "For decisions with tradeoffs, maps to ADR 'Alternatives Considered'"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Meta-prompting Pattern (Emerging)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Prompt generates prompts - create reusable templates"
    p.level = 1
    p.font.italic = True

    # Slide 4: Workflow Orchestration Strategies (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Workflow Orchestration Strategies"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Strategy A: ADR-Driven (Industry Standard)"

    p = tf.add_paragraph()
    p.text = "Write ADR → Reference in .cursorrules → Use AI → Validate"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Proven, lightweight, widely understood"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Strategy B: Structured File Workflow (Today's Example)"

    p = tf.add_paragraph()
    p.text = "System prompt → Task spec → Execution plan → Decisions → Generate"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Explicit pattern application, good for learning"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Strategy C: Tool-Assisted (Platform-Specific)"

    p = tf.add_paragraph()
    p.text = "Spec-Kit, GitHub Copilot Workspace, Cursor Composer"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Built-in guardrails, varying maturity levels"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "All strategies apply same underlying patterns!"
    p.font.bold = True
    p.font.size = Pt(20)

    # Slide 5: ReAct Pattern for Task Execution (3 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ReAct Pattern: Iterative Validated Steps"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = 'Without ReAct: "I\'ll migrate everything at once"'
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    p = tf.add_paragraph()
    p.text = "Something breaks, unclear which change caused it"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 0, 0)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "With ReAct: Iterative, Validated Steps"
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "THINK: What dependencies must be satisfied first?"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ACT: Replace javax.* with jakarta.* imports"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "OBSERVE: mvn compile → SUCCESS ✓"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "THINK: Now annotations can be safely updated..."
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Benefits: Logical ordering, early error detection, clear rollback"
    p.font.bold = True

    # Slide 6: Tree of Thoughts for Decision Making (2 min)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Tree of Thoughts: Exploring Alternatives"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "When to Use Tree of Thoughts:"

    p = tf.add_paragraph()
    p.text = "Multiple valid approaches exist"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Tradeoffs matter for your specific context"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Team needs to understand decision rationale"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Example: Exception Handling Strategy"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Branch A: Keep current (minimal changes, low risk) ✓"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Branch B: Adopt ProblemDetail (modern, more testing)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Branch C: Full RFC 7807 (over-engineered for scope)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Branch D: Simple strings (too basic)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Maps to ADR 'Alternatives Considered' section"
    p.font.bold = True
    p.font.size = Pt(18)

    # Save presentation
    output_file = "SESSION-2-Advanced-Patterns.pptx"
    prs.save(output_file)
    print(f"✅ Created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    print("🎨 Creating Session 2 PowerPoint slides (Advanced Patterns version)...")
    output = create_session2_slides_v2()
    print(f"\n🎉 Done! Open {output} to view slides")
